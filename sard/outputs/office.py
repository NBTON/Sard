"""GenOffice-inspired PowerPoint (.pptx) and Office document generator for Sard.

Provides automated authoring of cultural presentation decks with Sard's
cultural design palette (paper, ink, clay, date, olive, gold, card) and
native Arabic right-to-left layout conventions.
"""

from __future__ import annotations

import io
import logging
import os
import uuid
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger("sard.outputs.office")

# ---------------------------------------------------------------------------
# Sard Cultural Color Palette
# ---------------------------------------------------------------------------
COLOR_PAPER = RGBColor(0xF3, 0xEE, 0xE4)      # #F3EEE4
COLOR_PAPER_2 = RGBColor(0xE8, 0xE0, 0xD2)    # #E8E0D2
COLOR_INK = RGBColor(0x14, 0x12, 0x10)        # #141210
COLOR_CLAY = RGBColor(0xBE, 0x4A, 0x24)       # #BE4A24
COLOR_DATE = RGBColor(0x6E, 0x1F, 0x1F)       # #6E1F1F
COLOR_OLIVE = RGBColor(0x4A, 0x51, 0x3C)      # #4A513C
COLOR_GOLD = RGBColor(0xC4, 0xA4, 0x6A)       # #C4A46A
COLOR_CARD = RGBColor(0xFA, 0xF7, 0xF1)       # #FAF7F1
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_MUTED = RGBColor(0x8A, 0x81, 0x78)

FONT_HEADING = "Noto Naskh Arabic"
FONT_BODY = "IBM Plex Sans Arabic"
FONT_FALLBACK = "Arial"


@dataclass
class SlideCard:
    """Represents a structured card within a comparison or multi-column slide."""
    title: str
    subtitle: str = ""
    bullets: List[str] = field(default_factory=list)
    badge: str = ""
    accent_color: str = "clay"  # clay, gold, olive, date


@dataclass
class TimelineItem:
    """Represents a milestone in a chronological timeline slide."""
    year_or_era: str
    hijri_year: str = ""
    title: str = ""
    description: str = ""


@dataclass
class SlideContent:
    """Represents the semantic content of a single slide."""
    slide_type: str  # title, briefing, comparison, timeline, key_points, summary
    title: str
    subtitle: str = ""
    body_paragraphs: List[str] = field(default_factory=list)
    bullets: List[str] = field(default_factory=list)
    cards: List[SlideCard] = field(default_factory=list)
    timeline: List[TimelineItem] = field(default_factory=list)
    quote: str = ""
    quote_author: str = ""
    footer_text: str = "سرد — المستشار الثقافي للمملكة العربية السعودية"
    region_badge: str = ""


@dataclass
class PresentationDeck:
    """Represents an entire cultural briefing presentation."""
    title: str
    topic: str
    region: str = "المملكة العربية السعودية"
    author: str = "سرد (Sard Cultural Agent)"
    slides: List[SlideContent] = field(default_factory=list)
    deck_id: str = field(default_factory=lambda: f"deck-{uuid.uuid4().hex[:8]}")


class PresentationGenerator:
    """Generates PowerPoint (.pptx) decks adhering strictly to Sard cultural design tokens."""

    def __init__(self, default_output_dir: Optional[Path] = None):
        self.output_dir = default_output_dir or Path("output")
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def _get_accent_rgb(self, name: str) -> RGBColor:
        mapping = {
            "clay": COLOR_CLAY,
            "gold": COLOR_GOLD,
            "olive": COLOR_OLIVE,
            "date": COLOR_DATE,
            "ink": COLOR_INK,
        }
        return mapping.get(name.lower(), COLOR_CLAY)

    def build_pptx(self, deck: PresentationDeck) -> bytes:
        """Constructs the presentation in 16:9 widescreen layout and returns bytes."""
        prs = Presentation()
        # 16:9 widescreen dimensions (13.333 x 7.5 inches)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        blank_slide_layout = prs.slide_layouts[6]  # Completely blank layout

        for slide_data in deck.slides:
            slide = prs.slides.add_slide(blank_slide_layout)
            self._render_slide(prs, slide, slide_data, deck)

        stream = io.BytesIO()
        prs.save(stream)
        return stream.getvalue()

    def _set_background(self, slide, color: RGBColor = COLOR_PAPER):
        """Sets a full-slide solid background shape."""
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def _render_slide(self, prs, slide, data: SlideContent, deck: PresentationDeck):
        if data.slide_type == "title":
            self._set_background(slide, COLOR_INK)
            self._render_title_slide(slide, data, deck)
        elif data.slide_type == "comparison":
            self._set_background(slide, COLOR_PAPER)
            self._render_comparison_slide(slide, data)
        elif data.slide_type == "timeline":
            self._set_background(slide, COLOR_PAPER)
            self._render_timeline_slide(slide, data)
        elif data.slide_type == "summary":
            self._set_background(slide, COLOR_PAPER_2)
            self._render_summary_slide(slide, data)
        else:
            self._set_background(slide, COLOR_PAPER)
            self._render_briefing_slide(slide, data)

    def _render_title_slide(self, slide, data: SlideContent, deck: PresentationDeck):
        """Renders dark-ink title cover slide with gold & clay accents."""
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(0.8), Inches(11.333), Inches(0.08)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = COLOR_GOLD
        top_bar.line.fill.background()

        badge_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(0.6))
        tf_badge = badge_box.text_frame
        tf_badge.word_wrap = True
        p_b = tf_badge.paragraphs[0]
        p_b.alignment = PP_ALIGN.RIGHT
        run_b = p_b.add_run()
        run_b.text = f"✦ {deck.region} • توثيق ثقافي معتمد"
        run_b.font.name = FONT_BODY
        run_b.font.size = Pt(14)
        run_b.font.color.rgb = COLOR_GOLD
        run_b.font.bold = True

        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = data.title
        run.font.name = FONT_HEADING
        run.font.size = Pt(40)
        run.font.color.rgb = COLOR_PAPER
        run.font.bold = True

        if data.subtitle:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.RIGHT
            p2.space_before = Pt(14)
            run2 = p2.add_run()
            run2.text = data.subtitle
            run2.font.name = FONT_BODY
            run2.font.size = Pt(20)
            run2.font.color.rgb = COLOR_GOLD

        mid_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(8.333), Inches(5.0), Inches(4.0), Inches(0.05)
        )
        mid_bar.fill.solid()
        mid_bar.fill.fore_color.rgb = COLOR_CLAY
        mid_bar.line.fill.background()

        footer_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.333), Inches(0.8))
        tf_foot = footer_box.text_frame
        p_f = tf_foot.paragraphs[0]
        p_f.alignment = PP_ALIGN.RIGHT
        r_f = p_f.add_run()
        r_f.text = f"إعداد: {deck.author} | مبادرة التوثيق الثقافي"
        r_f.font.name = FONT_BODY
        r_f.font.size = Pt(13)
        r_f.font.color.rgb = COLOR_MUTED

    def _render_briefing_slide(self, slide, data: SlideContent):
        """Renders standard briefing slide with header bar and styled paragraphs."""
        self._add_slide_header(slide, data.title, data.region_badge)

        card_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8)
        )
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = COLOR_CARD
        card_bg.line.color.rgb = COLOR_PAPER_2
        card_bg.line.width = Pt(1)

        content_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(10.933), Inches(4.2))
        tf = content_box.text_frame
        tf.word_wrap = True

        first = True
        for p_text in data.body_paragraphs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.RIGHT
            p.space_after = Pt(12)
            run = p.add_run()
            run.text = p_text
            run.font.name = FONT_BODY
            run.font.size = Pt(16)
            run.font.color.rgb = COLOR_INK

        for bullet in data.bullets:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.RIGHT
            p.space_after = Pt(8)
            run_icon = p.add_run()
            run_icon.text = "✦  "
            run_icon.font.name = FONT_BODY
            run_icon.font.size = Pt(14)
            run_icon.font.color.rgb = COLOR_CLAY
            run_icon.font.bold = True

            run_b = p.add_run()
            run_b.text = bullet
            run_b.font.name = FONT_BODY
            run_b.font.size = Pt(15)
            run_b.font.color.rgb = COLOR_INK

        if data.quote:
            p_q = tf.add_paragraph()
            p_q.alignment = PP_ALIGN.RIGHT
            p_q.space_before = Pt(14)
            r_q = p_q.add_run()
            r_q.text = f"«{data.quote}»"
            r_q.font.name = FONT_HEADING
            r_q.font.size = Pt(16)
            r_q.font.color.rgb = COLOR_DATE
            r_q.font.italic = True

        self._add_slide_footer(slide, data.footer_text)

    def _render_comparison_slide(self, slide, data: SlideContent):
        """Renders multi-column comparison cards (e.g. Najdi vs Hejazi vs Asiri)."""
        self._add_slide_header(slide, data.title, data.region_badge or "مقارنة تراثية")

        cards = data.cards or []
        num_cards = max(1, min(len(cards), 3))
        gap = Inches(0.4)
        total_width = Inches(11.733)
        card_width = (total_width - (gap * (num_cards - 1))) / num_cards
        start_x = Inches(0.8)

        for idx, card in enumerate(cards[:3]):
            x = start_x + (idx * (card_width + gap))
            y = Inches(1.8)
            h = Inches(4.8)

            c_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, h)
            c_shape.fill.solid()
            c_shape.fill.fore_color.rgb = COLOR_CARD
            c_shape.line.color.rgb = COLOR_PAPER_2
            c_shape.line.width = Pt(1)

            accent_rgb = self._get_accent_rgb(card.accent_color)
            top_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, card_width, Inches(0.12)
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = accent_rgb
            top_bar.line.fill.background()

            tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.25), card_width - Inches(0.4), h - Inches(0.5))
            tf = tb.text_frame
            tf.word_wrap = True

            p_t = tf.paragraphs[0]
            p_t.alignment = PP_ALIGN.RIGHT
            r_t = p_t.add_run()
            r_t.text = card.title
            r_t.font.name = FONT_HEADING
            r_t.font.size = Pt(19)
            r_t.font.bold = True
            r_t.font.color.rgb = accent_rgb

            if card.subtitle:
                p_sub = tf.add_paragraph()
                p_sub.alignment = PP_ALIGN.RIGHT
                p_sub.space_after = Pt(8)
                r_sub = p_sub.add_run()
                r_sub.text = card.subtitle
                r_sub.font.name = FONT_BODY
                r_sub.font.size = Pt(13)
                r_sub.font.color.rgb = COLOR_MUTED

            for b in card.bullets:
                p_b = tf.add_paragraph()
                p_b.alignment = PP_ALIGN.RIGHT
                p_b.space_after = Pt(6)
                r_dot = p_b.add_run()
                r_dot.text = "• "
                r_dot.font.color.rgb = accent_rgb
                r_dot.font.bold = True
                r_text = p_b.add_run()
                r_text.text = b
                r_text.font.name = FONT_BODY
                r_text.font.size = Pt(14)
                r_text.font.color.rgb = COLOR_INK

        self._add_slide_footer(slide, data.footer_text)

    def _render_timeline_slide(self, slide, data: SlideContent):
        """Renders chronological timeline slide."""
        self._add_slide_header(slide, data.title, data.region_badge or "تسلسل تاريخي")

        timeline = data.timeline or []
        num_steps = max(1, min(len(timeline), 4))
        step_width = Inches(11.733) / num_steps
        start_x = Inches(0.8)

        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.4), Inches(11.733), Inches(0.06)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_GOLD
        line.line.fill.background()

        for idx, item in enumerate(timeline[:4]):
            x = start_x + (idx * step_width)

            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, x + (step_width / 2) - Inches(0.2), Inches(3.28), Inches(0.4), Inches(0.4)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLOR_CLAY
            circle.line.color.rgb = COLOR_PAPER
            circle.line.width = Pt(2)

            tb_year = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.2), step_width - Inches(0.2), Inches(1.0))
            tf_y = tb_year.text_frame
            tf_y.word_wrap = True
            p_y = tf_y.paragraphs[0]
            p_y.alignment = PP_ALIGN.CENTER
            r_y = p_y.add_run()
            r_y.text = item.year_or_era
            r_y.font.name = FONT_HEADING
            r_y.font.size = Pt(18)
            r_y.font.bold = True
            r_y.font.color.rgb = COLOR_DATE

            if item.hijri_year:
                p_h = tf_y.add_paragraph()
                p_h.alignment = PP_ALIGN.CENTER
                r_h = p_h.add_run()
                r_h.text = item.hijri_year
                r_h.font.name = FONT_BODY
                r_h.font.size = Pt(12)
                r_h.font.color.rgb = COLOR_GOLD

            tb_desc = slide.shapes.add_textbox(x + Inches(0.1), Inches(3.9), step_width - Inches(0.2), Inches(2.7))
            tf_d = tb_desc.text_frame
            tf_d.word_wrap = True
            p_dt = tf_d.paragraphs[0]
            p_dt.alignment = PP_ALIGN.CENTER
            r_dt = p_dt.add_run()
            r_dt.text = item.title
            r_dt.font.name = FONT_BODY
            r_dt.font.size = Pt(15)
            r_dt.font.bold = True
            r_dt.font.color.rgb = COLOR_INK

            if item.description:
                p_dd = tf_d.add_paragraph()
                p_dd.alignment = PP_ALIGN.CENTER
                p_dd.space_before = Pt(6)
                r_dd = p_dd.add_run()
                r_dd.text = item.description
                r_dd.font.name = FONT_BODY
                r_dd.font.size = Pt(13)
                r_dd.font.color.rgb = COLOR_INK

        self._add_slide_footer(slide, data.footer_text)

    def _render_summary_slide(self, slide, data: SlideContent):
        """Renders closing / takeaways slide with highlighted takeaway card."""
        self._add_slide_header(slide, data.title, "خلاصة وتوثيق")

        card_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.8), Inches(10.333), Inches(4.8)
        )
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = COLOR_CARD
        card_bg.line.color.rgb = COLOR_GOLD
        card_bg.line.width = Pt(1.5)

        tb = slide.shapes.add_textbox(Inches(1.9), Inches(2.1), Inches(9.533), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True

        first = True
        for b in data.bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.RIGHT
            p.space_after = Pt(12)
            r_i = p.add_run()
            r_i.text = "✦  "
            r_i.font.color.rgb = COLOR_CLAY
            r_i.font.bold = True
            r_i.font.size = Pt(16)

            r_t = p.add_run()
            r_t.text = b
            r_t.font.name = FONT_BODY
            r_t.font.size = Pt(16)
            r_t.font.color.rgb = COLOR_INK

        if data.quote:
            p_q = tf.add_paragraph()
            p_q.alignment = PP_ALIGN.CENTER
            p_q.space_before = Pt(16)
            r_q = p_q.add_run()
            r_q.text = f"«{data.quote}»"
            r_q.font.name = FONT_HEADING
            r_q.font.size = Pt(17)
            r_q.font.color.rgb = COLOR_DATE
            r_q.font.italic = True

        self._add_slide_footer(slide, data.footer_text)

    def _add_slide_header(self, slide, title: str, badge: str = ""):
        """Renders standard top bar with title and cultural badge."""
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.04)
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = COLOR_CLAY
        stripe.line.fill.background()

        h_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(8.5), Inches(0.9))
        tf = h_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = title
        r.font.name = FONT_HEADING
        r.font.size = Pt(24)
        r.font.bold = True
        r.font.color.rgb = COLOR_INK

        if badge:
            b_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.65), Inches(3.0), Inches(0.6))
            tf_b = b_box.text_frame
            p_b = tf_b.paragraphs[0]
            p_b.alignment = PP_ALIGN.LEFT
            r_b = p_b.add_run()
            r_b.text = f"[{badge}]"
            r_b.font.name = FONT_BODY
            r_b.font.size = Pt(13)
            r_b.font.color.rgb = COLOR_OLIVE
            r_b.font.bold = True

    def _add_slide_footer(self, slide, footer_text: str):
        """Renders discrete footer with branding."""
        f_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.733), Inches(0.4))
        tf = f_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = footer_text
        r.font.name = FONT_BODY
        r.font.size = Pt(10)
        r.font.color.rgb = COLOR_MUTED

    def save_deck_file(self, deck: PresentationDeck, filename: Optional[str] = None) -> Tuple[Path, str]:
        """Builds and writes presentation to a .pptx file."""
        safe_name = filename or f"{deck.deck_id}.pptx"
        if not safe_name.endswith(".pptx"):
            safe_name += ".pptx"
        target_path = self.output_dir / safe_name
        data = self.build_pptx(deck)
        target_path.write_bytes(data)
        logger.info("Saved PowerPoint presentation deck: %s (%d bytes)", target_path, len(data))
        return target_path, safe_name


# ---------------------------------------------------------------------------
# High-Level Templates & Presets
# ---------------------------------------------------------------------------

def create_cultural_briefing_deck(
    topic: str,
    region: str = "المملكة العربية السعودية",
    overview_text: str = "",
    comparison_cards: Optional[List[Dict[str, Any]]] = None,
    timeline_items: Optional[List[Dict[str, Any]]] = None,
    key_takeaways: Optional[List[str]] = None,
    quote: str = "",
) -> PresentationDeck:
    """Creates a complete multi-slide cultural briefing presentation deck."""
    deck = PresentationDeck(
        title=f"الإيجاز الثقافي: {topic}",
        topic=topic,
        region=region,
    )

    # 1. Title Slide
    deck.slides.append(
        SlideContent(
            slide_type="title",
            title=topic,
            subtitle=f"إيجاز توثيقي شامل حول التراث والأصالة في {region}",
            region_badge=region,
        )
    )

    # 2. Overview Slide
    overview_bullets = [
        f"استعراض شامل للجذور التاريخية والمعالم الثقافية في {region}.",
        "توثيق العادات والتقاليد المتوارثة عبر الأجيال وصيانتها.",
        "التوافق مع أحدث معايير هيئة التراث ووزارة الثقافة السعودية.",
    ]
    deck.slides.append(
        SlideContent(
            slide_type="briefing",
            title=f"مدخل إلى {topic}",
            body_paragraphs=[overview_text] if overview_text else [
                f"يمثل «{topic}» ركيزة أساسية من ركائز الهوية الثقافية السعودية التي تعكس عمق التاريخ وتنوع البيئات والمناطق."
            ],
            bullets=overview_bullets,
            quote=quote,
            region_badge=region,
        )
    )

    # 3. Comparative Analysis Slide (if provided)
    if comparison_cards:
        cards_objs = []
        for c in comparison_cards:
            cards_objs.append(
                SlideCard(
                    title=c.get("title", ""),
                    subtitle=c.get("subtitle", ""),
                    bullets=c.get("bullets", []),
                    badge=c.get("badge", ""),
                    accent_color=c.get("accent_color", "clay"),
                )
            )
        deck.slides.append(
            SlideContent(
                slide_type="comparison",
                title=f"الخصائص والأنماط الإقليمية: {topic}",
                cards=cards_objs,
                region_badge=region,
            )
        )

    # 4. Timeline Slide (if provided)
    if timeline_items:
        t_objs = []
        for t in timeline_items:
            t_objs.append(
                TimelineItem(
                    year_or_era=t.get("year_or_era", ""),
                    hijri_year=t.get("hijri_year", ""),
                    title=t.get("title", ""),
                    description=t.get("description", ""),
                )
            )
        deck.slides.append(
            SlideContent(
                slide_type="timeline",
                title=f"المحطات التاريخية والتحولات: {topic}",
                timeline=t_objs,
                region_badge=region,
            )
        )

    # 5. Summary / Takeaways Slide
    summary_bullets = key_takeaways or [
        f"أهمية صون عناصر {topic} ونقلها للأجيال القادمة كرمز للأصالة.",
        "الارتباط الوثيق بين البيئة الطبيعية والابتكار الإنساني في المملكة.",
        "الاستناد إلى المراجع المعتمدة (دارة الملك عبد العزيز، هيئة التراث، وزارة الثقافة).",
    ]
    deck.slides.append(
        SlideContent(
            slide_type="summary",
            title="الخلاصة والتوصيات التراثية",
            bullets=summary_bullets,
            quote="تراثنا هويتنا، وأصالتنا جسر نحو المستقبل.",
            region_badge=region,
        )
    )

    return deck
